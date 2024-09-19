#!/usr/bin/python3
"""
file: pptx2audio.py
description: Convert your pptx to audiobook in M4B format
Usage example:
    python pptx2audio.py presentation.pptx
"""
import sys
import os
import logging
import asyncio
import pptx
from pptx import presentation, slide
from backend_audio import m4b
from frontend import input_tool

logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)

if sys.platform in ("win32", "cygwin"):
    BACK_END_TTS = "EDGE_TTS"
elif sys.platform == "darwin":
    BACK_END_TTS = '"GTTS'
else:
    BACK_END_TTS = "PYTTS"
LANGUAGE_DICT = {"it-IT":"it"}
LANGUAGE_DICT_PYTTS = {"it":"italian", "en":"english"}

def __get_notes(note: slide.NotesSlide) -> str:
    print(note)
    return ""
def __extract_text_from_slide(slide_obj: slide.Slide) -> str:
    text_slide, text_note = "", ""
    if slide_obj.has_notes_slide:
        text_note = __get_notes(slide_obj.notes_slide)
    for shape in slide_obj.shapes:
        if shape.has_text_frame:
            text_slide += shape.text
    text_slide = text_slide.strip()
    if len(text_note) != 0:
        text_slide += f"\n\nNote:\n{text_note}"
    return text_slide

def extract_pptx_text(path_pptx:str) -> str:
    """Extract the text from each slide and concat.
    Args:
        path_pptx (str): pptx presentation path
    Return:
        str: all text generated by presentation
    """
    text_out:str = ""
    p: presentation.Presentation = pptx.Presentation(path_pptx)
    for idx,s in enumerate(p.slides):
        text_out += f"\n\nSlide numero: {idx}\n"
        ts = __extract_text_from_slide(s)
        text_out += ts
        if len(ts) == 0:
            text_out += ""
    return text_out

if __name__ == "__main__":
    input_file_path, output_file_path = input_tool.get_sys_input(os.path.dirname(__file__))
    chapters = []
    metada_output = {}
    ch_metadatas = []

    m4b.init(BACK_END_TTS)

    text=extract_pptx_text(input_file_path)
    if BACK_END_TTS == "PYTTS":
        m4b.generate_audio_pytts(text, output_file_path, lang="en")
    else:
        loop = asyncio.get_event_loop_policy().get_event_loop()
        try:
            voices = loop.run_until_complete(m4b.get_voices_edge_tts(lang="it"))
            VOICE = voices[0]["Name"]
            loop.run_until_complete(m4b.generate_audio_edge_tts(text, output_file_path,
                                                                lang="it-IT", voice=VOICE))
        finally:
            loop.close()
    #generate_m4b(output_file_path, chapters, metada_output, chapter_metadata=ch_metadatas)
